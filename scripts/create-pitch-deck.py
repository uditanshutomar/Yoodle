#!/usr/bin/env python3
"""
Yoodle HackCU Pitch Deck Generator
Neo-brutalist theme with bold colors, thick borders, and playful design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme Colors ──────────────────────────────────────────────────────
BG_CREAM    = RGBColor(0xFA, 0xFA, 0xF8)
BG_DARK     = RGBColor(0x0A, 0x0A, 0x0A)
YELLOW      = RGBColor(0xFF, 0xE6, 0x00)
CORAL       = RGBColor(0xFF, 0x6B, 0x6B)
VIOLET      = RGBColor(0x7C, 0x3A, 0xED)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x0A, 0x0A, 0x0A)
GRAY        = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xEE)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=18, default_color=BLACK, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add text box with multiple styled paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.get("text", "")
        p.font.size = Pt(line.get("size", default_size))
        p.font.color.rgb = line.get("color", default_color)
        p.font.bold = line.get("bold", False)
        p.font.name = line.get("font", font_name)
        p.alignment = line.get("align", align)
        p.space_after = Pt(line.get("space_after", 6))
        p.space_before = Pt(line.get("space_before", 0))
    return txBox

def add_shadow_rect(slide, left, top, width, height, fill_color, shadow_offset=Inches(0.08)):
    """Add a card with neo-brutalist shadow effect."""
    # Shadow (offset dark rect)
    add_rect(slide, left + shadow_offset, top + shadow_offset, width, height, BLACK)
    # Main card
    card = add_rect(slide, left, top, width, height, fill_color, BLACK, Pt(3))
    return card

def add_shadow_rounded(slide, left, top, width, height, fill_color, shadow_offset=Inches(0.08)):
    """Add a rounded card with neo-brutalist shadow."""
    add_rounded_rect(slide, left + shadow_offset, top + shadow_offset, width, height, BLACK)
    card = add_rounded_rect(slide, left, top, width, height, fill_color, BLACK, Pt(3))
    return card

def add_circle(slide, left, top, size, fill_color, border_color=None, border_width=Pt(0)):
    """Add a circle/oval."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

# ── SLIDE 1: Title ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_CREAM)

# Yellow accent bar at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), YELLOW)

# Decorative circles
add_circle(slide, Inches(10.5), Inches(0.5), Inches(1.2), YELLOW, BLACK, Pt(3))
add_circle(slide, Inches(11.2), Inches(1.3), Inches(0.6), CORAL, BLACK, Pt(2))
add_circle(slide, Inches(0.5), Inches(5.5), Inches(0.8), VIOLET, BLACK, Pt(2))
add_circle(slide, Inches(1.5), Inches(6.0), Inches(0.5), CYAN, BLACK, Pt(2))

# Logo image
logo_path = "/Users/uditanshutomar/Desktop/Yoodle/yoodle-logo.png"
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(1.5), Inches(1.0), Inches(2.5), Inches(2.5))

# Title
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(1.2),
             "YOODLE", font_size=72, color=BLACK, bold=True, font_name="Arial Black")

# Tagline
add_text_box(slide, Inches(1.5), Inches(4.7), Inches(8), Inches(0.6),
             "Meetings that actually slap. \U0001f91f", font_size=32, color=BLACK, bold=False, font_name="Arial")

# Subtitle
add_text_box(slide, Inches(1.5), Inches(5.4), Inches(8), Inches(0.5),
             "Video calls. AI notes. Shared workspaces. Built for how you actually work.",
             font_size=18, color=GRAY, font_name="Arial")

# HackCU badge
add_shadow_rounded(slide, Inches(1.5), Inches(6.2), Inches(3.2), Inches(0.7), YELLOW)
add_text_box(slide, Inches(1.6), Inches(6.28), Inches(3), Inches(0.55),
             "\U0001f3d7\ufe0f  HackCU 2025", font_size=22, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")


# ── SLIDE 2: The Problem ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# Yellow accent
add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, YELLOW)

# Section tag
add_shadow_rounded(slide, Inches(1.0), Inches(0.6), Inches(3.5), Inches(0.65), CORAL)
add_text_box(slide, Inches(1.1), Inches(0.67), Inches(3.3), Inches(0.5),
             "\U0001f6a8  THE PROBLEM", font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

# Main heading
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(1.0),
             "Meetings are broken.", font_size=54, color=WHITE, bold=True, font_name="Arial Black")

# Problem points
problems = [
    {"text": "\u274c  Zoom fatigue is real — endless, draining video calls", "size": 24, "color": WHITE, "bold": False, "space_after": 14},
    {"text": "\u274c  Notes get lost — nobody remembers action items", "size": 24, "color": WHITE, "bold": False, "space_after": 14},
    {"text": "\u274c  Collaboration is fragmented — code on Replit, chat on Discord, calls on Zoom", "size": 24, "color": WHITE, "bold": False, "space_after": 14},
    {"text": "\u274c  Sensitive convos leave digital footprints forever", "size": 24, "color": WHITE, "bold": False, "space_after": 14},
    {"text": "\u274c  No tool is built for how Gen Z actually works", "size": 24, "color": WHITE, "bold": False, "space_after": 14},
]
add_multi_text(slide, Inches(1.0), Inches(3.0), Inches(11), Inches(3.5), problems)

# Stat callout
add_shadow_rounded(slide, Inches(8.5), Inches(5.8), Inches(4.0), Inches(1.2), YELLOW)
add_multi_text(slide, Inches(8.7), Inches(5.9), Inches(3.6), Inches(1.0), [
    {"text": "73% of Gen Z", "size": 28, "color": BLACK, "bold": True, "space_after": 2},
    {"text": "say current tools don't fit their workflow", "size": 16, "color": BLACK, "bold": False},
])


# ── SLIDE 3: Our Solution ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_CREAM)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), YELLOW)

# Section tag
add_shadow_rounded(slide, Inches(1.0), Inches(0.6), Inches(3.5), Inches(0.65), YELLOW)
add_text_box(slide, Inches(1.1), Inches(0.67), Inches(3.3), Inches(0.5),
             "\U0001f4a1  OUR SOLUTION", font_size=22, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

# Main heading
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(11), Inches(0.8),
             "One platform. Everything you need.", font_size=44, color=BLACK, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(1.0), Inches(2.6), Inches(10), Inches(0.5),
             "Yoodle combines video calling, AI intelligence, shared workspaces, and ephemeral rooms into one seamless experience.",
             font_size=20, color=GRAY, font_name="Arial")

# 4 Feature cards
card_data = [
    {"title": "\U0001f4f9 Crystal Calls", "desc": "HD video & audio with WebRTC peer-to-peer connections", "color": YELLOW},
    {"title": "\U0001f9e0 AI Brain", "desc": "Auto transcripts, meeting summaries, action items", "color": VIOLET},
    {"title": "\U0001f47b Ghost Rooms", "desc": "Ephemeral spaces that vanish — nothing is stored", "color": CORAL},
    {"title": "\U0001f680 Ship Together", "desc": "Shared cloud VMs with terminals & collaborative coding", "color": CYAN},
]

for i, card in enumerate(card_data):
    x = Inches(1.0 + i * 3.0)
    y = Inches(3.6)
    w = Inches(2.7)
    h = Inches(3.2)
    add_shadow_rounded(slide, x, y, w, h, WHITE)
    # Colored top bar on card
    add_rect(slide, x + Inches(0.02), y + Inches(0.02), w - Inches(0.04), Inches(0.4), card["color"])
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), w - Inches(0.4), Inches(0.6),
                 card["title"], font_size=22, color=BLACK, bold=True, font_name="Arial Black")
    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), w - Inches(0.4), Inches(1.5),
                 card["desc"], font_size=16, color=GRAY, font_name="Arial")


# ── SLIDE 4: Crystal Calls ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_CREAM)

# Yellow left accent
add_rect(slide, Inches(0), Inches(0), Inches(6.0), SLIDE_H, YELLOW)
add_rect(slide, Inches(6.0), Inches(0), Inches(0.06), SLIDE_H, BLACK)

# Left side content
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(4.5), Inches(0.6),
             "\U0001f4f9 CRYSTAL CALLS", font_size=20, color=BLACK, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(4.8), Inches(1.0),
             "Video calls that don't suck.", font_size=42, color=BLACK, bold=True, font_name="Arial Black")

features_cc = [
    {"text": "\u2022  WebRTC peer-to-peer — zero lag, HD quality", "size": 20, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Screen sharing with one click", "size": 20, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Voice activity detection — see who's talking", "size": 20, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Mixed-audio recording (all participants)", "size": 20, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  TURN/STUN server for NAT traversal", "size": 20, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Real-time media state sync via Socket.io", "size": 20, "color": BLACK, "bold": False, "space_after": 12},
]
add_multi_text(slide, Inches(0.8), Inches(3.0), Inches(4.8), Inches(3.5), features_cc)

# Right side — decorative video grid mockup
add_shadow_rounded(slide, Inches(7.0), Inches(1.0), Inches(5.0), Inches(5.0), WHITE)

# Mock video tiles (circles like Yoodle's design)
circle_positions = [
    (Inches(8.0), Inches(1.8), Inches(1.6), CORAL),
    (Inches(10.0), Inches(1.8), Inches(1.6), CYAN),
    (Inches(8.0), Inches(3.8), Inches(1.6), VIOLET),
    (Inches(10.0), Inches(3.8), Inches(1.6), YELLOW),
]
for (cx, cy, cs, cc) in circle_positions:
    add_circle(slide, cx, cy, cs, cc, BLACK, Pt(3))

add_text_box(slide, Inches(7.5), Inches(5.6), Inches(4.0), Inches(0.4),
             "Floating circles — not boring boxes", font_size=14, color=GRAY, align=PP_ALIGN.CENTER, font_name="Arial")


# ── SLIDE 5: AI Brain ────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), VIOLET)

# Section tag
add_shadow_rounded(slide, Inches(1.0), Inches(0.6), Inches(3.0), Inches(0.65), VIOLET)
add_text_box(slide, Inches(1.1), Inches(0.67), Inches(2.8), Inches(0.5),
             "\U0001f9e0  AI BRAIN", font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text_box(slide, Inches(1.0), Inches(1.7), Inches(11), Inches(0.8),
             "Your meetings, remembered.", font_size=48, color=WHITE, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(1.0), Inches(2.6), Inches(10), Inches(0.5),
             "Every meeting auto-generates transcripts, summaries, and action items — powered by Gemini & ElevenLabs.",
             font_size=20, color=RGBColor(0xBB, 0xBB, 0xBB), font_name="Arial")

# 3 AI feature cards
ai_cards = [
    {"icon": "\U0001f399\ufe0f", "title": "Live Transcription", "desc": "3-second audio chunks sent to ElevenLabs STT. Real-time captions for all participants.", "color": VIOLET},
    {"icon": "\U0001f4dd", "title": "Smart Summaries", "desc": "Gemini generates meeting notes with key decisions, action items, and deadlines.", "color": CORAL},
    {"icon": "\U0001f4be", "title": "Recording + Playback", "desc": "Mixed-audio recordings uploaded to Vultr S3. Download anytime from the meeting page.", "color": CYAN},
]

for i, card in enumerate(ai_cards):
    x = Inches(1.0 + i * 3.8)
    y = Inches(3.6)
    w = Inches(3.4)
    h = Inches(3.2)
    add_shadow_rounded(slide, x, y, w, h, RGBColor(0x1A, 0x1A, 0x1A))
    # Colored accent bar
    add_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.08), card["color"])
    add_text_box(slide, x + Inches(0.2), y + Inches(0.4), w - Inches(0.4), Inches(0.5),
                 f'{card["icon"]}  {card["title"]}', font_size=22, color=WHITE, bold=True, font_name="Arial Black")
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), w - Inches(0.4), Inches(1.8),
                 card["desc"], font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Arial")


# ── SLIDE 6: Ghost Rooms ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_CREAM)

# Coral right accent
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_CREAM)
add_rect(slide, Inches(7.5), Inches(0), Inches(5.833), SLIDE_H, CORAL)
add_rect(slide, Inches(7.5), Inches(0), Inches(0.06), SLIDE_H, BLACK)

# Left content
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(6.0), Inches(0.6),
             "\U0001f47b GHOST ROOMS", font_size=20, color=BLACK, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6.2), Inches(1.0),
             "Talk freely.\nNothing stays.", font_size=44, color=BLACK, bold=True, font_name="Arial Black")

ghost_features = [
    {"text": "\u2022  Ephemeral rooms — auto-delete after 4 hours (MongoDB TTL)", "size": 19, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Real-time code sharing with syntax highlighting", "size": 19, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Messages, notes, participants — all temporary", "size": 19, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Democratic save: data persists ONLY if all vote yes", "size": 19, "color": BLACK, "bold": False, "space_after": 12},
    {"text": "\u2022  Perfect for brainstorms, sensitive convos, quick syncs", "size": 19, "color": BLACK, "bold": False, "space_after": 12},
]
add_multi_text(slide, Inches(0.8), Inches(3.2), Inches(6.2), Inches(3.5), ghost_features)

# Right side — ghost illustration
add_text_box(slide, Inches(8.5), Inches(2.5), Inches(3.5), Inches(2.0),
             "\U0001f47b", font_size=120, color=WHITE, bold=False, align=PP_ALIGN.CENTER, font_name="Arial")
add_text_box(slide, Inches(8.5), Inches(4.5), Inches(3.5), Inches(0.5),
             "poof — it's gone", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")
add_text_box(slide, Inches(8.2), Inches(5.2), Inches(4.0), Inches(0.8),
             "Surviving Vercel cold starts with MongoDB-backed ephemeral storage + TTL indexes",
             font_size=14, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER, font_name="Arial")


# ── SLIDE 7: Ship Together (Workspaces) ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), CYAN)

add_shadow_rounded(slide, Inches(1.0), Inches(0.6), Inches(4.0), Inches(0.65), CYAN)
add_text_box(slide, Inches(1.1), Inches(0.67), Inches(3.8), Inches(0.5),
             "\U0001f680  SHIP TOGETHER", font_size=22, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text_box(slide, Inches(1.0), Inches(1.7), Inches(11), Inches(0.8),
             "Cloud VMs for your whole team.", font_size=48, color=WHITE, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(1.0), Inches(2.6), Inches(10), Inches(0.5),
             "Shared development environments powered by Vultr. One VM, everyone codes together.",
             font_size=20, color=RGBColor(0xBB, 0xBB, 0xBB), font_name="Arial")

# Feature cards in 2x2 grid
ws_cards = [
    {"icon": "\U0001f4bb", "title": "Web Terminal", "desc": "Full xterm.js terminal in the browser via SSH proxy", "x": 1.0, "y": 3.5},
    {"icon": "\u2601\ufe0f", "title": "Vultr VMs", "desc": "1-click provisioning with cloud-init (Node, Docker, Git)", "x": 6.8, "y": 3.5},
    {"icon": "\U0001f465", "title": "Team Access", "desc": "All workspace members share the same machine", "x": 1.0, "y": 5.3},
    {"icon": "\U0001f510", "title": "SSH Key Auth", "desc": "Secure access with UUID-validated SSH keys + password fallback", "x": 6.8, "y": 5.3},
]

for card in ws_cards:
    x = Inches(card["x"])
    y = Inches(card["y"])
    w = Inches(5.2)
    h = Inches(1.4)
    add_shadow_rounded(slide, x, y, w, h, RGBColor(0x1A, 0x1A, 0x1A))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
                 f'{card["icon"]}  {card["title"]}', font_size=20, color=CYAN, bold=True, font_name="Arial Black")
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(0.6),
                 card["desc"], font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Arial")


# ── SLIDE 8: Doodle Poodle (AI Assistant) ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_CREAM)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), VIOLET)

# Left side
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(6.0), Inches(0.6),
             "\U0001f9e0 DOODLE POODLE", font_size=20, color=BLACK, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6.2), Inches(1.0),
             "Your AI sidekick\nthat actually helps.", font_size=42, color=BLACK, bold=True, font_name="Arial Black")

poodle_features = [
    {"text": "\u2022  Powered by Google Gemini API", "size": 20, "color": BLACK, "bold": False, "space_after": 10},
    {"text": "\u2022  Meeting prep — context on what you need to know", "size": 20, "color": BLACK, "bold": False, "space_after": 10},
    {"text": "\u2022  Summarizes plans, proofreads docs", "size": 20, "color": BLACK, "bold": False, "space_after": 10},
    {"text": "\u2022  Remembers details across conversations", "size": 20, "color": BLACK, "bold": False, "space_after": 10},
    {"text": "\u2022  Chat interface with real-time streaming", "size": 20, "color": BLACK, "bold": False, "space_after": 10},
]
add_multi_text(slide, Inches(0.8), Inches(3.2), Inches(6.0), Inches(3.5), poodle_features)

# Right side — AI chat mockup
add_shadow_rounded(slide, Inches(7.5), Inches(0.8), Inches(5.0), Inches(6.0), WHITE)
add_rect(slide, Inches(7.5), Inches(0.8), Inches(5.0), Inches(0.7), VIOLET)
add_text_box(slide, Inches(7.8), Inches(0.88), Inches(4.4), Inches(0.5),
             "\U0001f9e0 Doodle Poodle", font_size=18, color=WHITE, bold=True, font_name="Arial Black")

# Chat bubbles
add_rounded_rect(slide, Inches(9.5), Inches(1.9), Inches(2.7), Inches(0.7), LIGHT_GRAY, BLACK, Pt(2))
add_text_box(slide, Inches(9.6), Inches(1.95), Inches(2.5), Inches(0.6),
             "Summarize my last meeting", font_size=14, color=BLACK, font_name="Arial")

add_rounded_rect(slide, Inches(7.8), Inches(2.9), Inches(4.0), Inches(1.6), VIOLET)
add_text_box(slide, Inches(7.9), Inches(2.95), Inches(3.8), Inches(1.5),
             "Your standup had 3 action items:\n1. Fix auth bug (assigned: you)\n2. Deploy v2 by Friday\n3. Review Sarah's PR",
             font_size=14, color=WHITE, font_name="Arial")

add_rounded_rect(slide, Inches(9.5), Inches(4.8), Inches(2.7), Inches(0.7), LIGHT_GRAY, BLACK, Pt(2))
add_text_box(slide, Inches(9.6), Inches(4.85), Inches(2.5), Inches(0.6),
             "What's on my plate today?", font_size=14, color=BLACK, font_name="Arial")

add_rounded_rect(slide, Inches(7.8), Inches(5.8), Inches(4.0), Inches(0.8), VIOLET)
add_text_box(slide, Inches(7.9), Inches(5.85), Inches(3.8), Inches(0.7),
             "You have the auth bug to fix and a\nmeeting with the design team at 3pm.",
             font_size=14, color=WHITE, font_name="Arial")


# ── SLIDE 9: Tech Stack ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), YELLOW)

add_text_box(slide, Inches(1.0), Inches(0.6), Inches(11), Inches(0.8),
             "\u2699\ufe0f  TECH STACK", font_size=20, color=YELLOW, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(1.0), Inches(1.3), Inches(11), Inches(0.8),
             "Built with sponsor tech + modern tools.", font_size=44, color=WHITE, bold=True, font_name="Arial Black")

# Tech categories
tech_groups = [
    {
        "title": "\U0001f3af Sponsor Tech",
        "items": "Gemini API \u2022 ElevenLabs \u2022 MongoDB Atlas \u2022 Vultr Cloud",
        "color": YELLOW,
        "x": 1.0, "y": 2.6, "w": 5.5
    },
    {
        "title": "\U0001f527 Frontend",
        "items": "Next.js 15 \u2022 React 19 \u2022 TypeScript 5 \u2022 Tailwind CSS 4 \u2022 Framer Motion",
        "color": CYAN,
        "x": 7.0, "y": 2.6, "w": 5.5
    },
    {
        "title": "\u26a1 Backend",
        "items": "Next.js API Routes \u2022 Mongoose 9 \u2022 Socket.io \u2022 WebRTC \u2022 jose (JWT)",
        "color": VIOLET,
        "x": 1.0, "y": 4.3, "w": 5.5
    },
    {
        "title": "\u2601\ufe0f Infrastructure",
        "items": "Vercel (app) \u2022 Vultr VPS (signaling + TURN) \u2022 Vultr S3 (recordings) \u2022 MongoDB Atlas",
        "color": CORAL,
        "x": 7.0, "y": 4.3, "w": 5.5
    },
    {
        "title": "\U0001f3a4 Real-time",
        "items": "WebRTC P2P \u2022 Socket.io signaling \u2022 MediaRecorder API \u2022 Web Audio API mixing \u2022 xterm.js",
        "color": YELLOW,
        "x": 1.0, "y": 6.0, "w": 11.5
    },
]

for tg in tech_groups:
    x = Inches(tg["x"])
    y = Inches(tg["y"])
    w = Inches(tg["w"])
    h = Inches(1.3)
    add_shadow_rounded(slide, x, y, w, h, RGBColor(0x1A, 0x1A, 0x1A))
    add_rect(slide, x + Inches(0.15), y + Inches(0.12), Inches(0.5), Inches(0.08), tg["color"])
    add_text_box(slide, x + Inches(0.2), y + Inches(0.25), w - Inches(0.4), Inches(0.4),
                 tg["title"], font_size=18, color=tg["color"], bold=True, font_name="Arial Black")
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(0.5),
                 tg["items"], font_size=15, color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Arial")


# ── SLIDE 10: Architecture ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_CREAM)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), YELLOW)

add_text_box(slide, Inches(1.0), Inches(0.6), Inches(11), Inches(0.6),
             "\U0001f3d7\ufe0f  ARCHITECTURE", font_size=20, color=BLACK, bold=True, font_name="Arial Black")

add_text_box(slide, Inches(1.0), Inches(1.2), Inches(11), Inches(0.8),
             "How it all connects.", font_size=40, color=BLACK, bold=True, font_name="Arial Black")

# Architecture flow boxes
arch_boxes = [
    {"label": "Browser\n(React 19)", "color": YELLOW, "x": 0.5, "y": 3.2, "w": 2.0, "h": 1.6},
    {"label": "Vercel\n(Next.js API)", "color": CYAN, "x": 3.3, "y": 3.2, "w": 2.2, "h": 1.6},
    {"label": "MongoDB\nAtlas", "color": VIOLET, "x": 6.3, "y": 3.2, "w": 2.0, "h": 1.6},
    {"label": "Vultr VPS\n(Signal+TURN)", "color": CORAL, "x": 3.3, "y": 5.5, "w": 2.2, "h": 1.4},
    {"label": "Vultr S3\n(Recordings)", "color": YELLOW, "x": 6.3, "y": 5.5, "w": 2.0, "h": 1.4},
    {"label": "Vultr VM\n(Workspace)", "color": CYAN, "x": 9.2, "y": 3.2, "w": 2.0, "h": 1.6},
    {"label": "Gemini +\nElevenLabs", "color": CORAL, "x": 9.2, "y": 5.5, "w": 2.0, "h": 1.4},
]

for ab in arch_boxes:
    x = Inches(ab["x"])
    y = Inches(ab["y"])
    w = Inches(ab["w"])
    h = Inches(ab["h"])
    add_shadow_rounded(slide, x, y, w, h, ab["color"])
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), h - Inches(0.3),
                 ab["label"], font_size=16, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

# Arrows (using text arrows for simplicity)
add_text_box(slide, Inches(2.5), Inches(3.7), Inches(0.8), Inches(0.5),
             "\u27a1", font_size=28, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.5), Inches(3.7), Inches(0.8), Inches(0.5),
             "\u27a1", font_size=28, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.3), Inches(3.7), Inches(0.9), Inches(0.5),
             "\u27a1", font_size=28, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.2), Inches(4.9), Inches(0.5), Inches(0.5),
             "\u2b07", font_size=28, color=BLACK, bold=True, align=PP_ALIGN.CENTER)


# ── SLIDE 11: Demo ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, YELLOW)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), BLACK)

# Decorative
add_circle(slide, Inches(0.5), Inches(1.0), Inches(1.5), CORAL, BLACK, Pt(3))
add_circle(slide, Inches(10.5), Inches(5.0), Inches(1.8), VIOLET, BLACK, Pt(3))
add_circle(slide, Inches(11.5), Inches(1.0), Inches(1.0), CYAN, BLACK, Pt(3))

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "DEMO TIME \U0001f3ac", font_size=72, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text_box(slide, Inches(2.0), Inches(3.5), Inches(9), Inches(0.6),
             "Let's see Yoodle in action.", font_size=32, color=BLACK, bold=False, align=PP_ALIGN.CENTER, font_name="Arial")

add_shadow_rounded(slide, Inches(4.0), Inches(4.8), Inches(5.3), Inches(0.8), WHITE)
add_text_box(slide, Inches(4.2), Inches(4.88), Inches(5.0), Inches(0.6),
             "\U0001f310  yoodle.vercel.app", font_size=26, color=BLACK, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text_box(slide, Inches(2.0), Inches(6.0), Inches(9), Inches(0.5),
             "Create a meeting \u2022 Join with a friend \u2022 AI transcription \u2022 Ghost Room \u2022 Terminal in VM",
             font_size=16, color=GRAY, align=PP_ALIGN.CENTER, font_name="Arial")


# ── SLIDE 12: Thank You ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), YELLOW)
add_rect(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), YELLOW)

# Decorative
add_circle(slide, Inches(0.5), Inches(0.5), Inches(1.0), YELLOW, BLACK, Pt(3))
add_circle(slide, Inches(11.0), Inches(5.5), Inches(1.5), CORAL, BLACK, Pt(3))
add_circle(slide, Inches(10.0), Inches(0.5), Inches(0.8), VIOLET, BLACK, Pt(3))

# Logo
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(5.4), Inches(0.8), Inches(2.5), Inches(2.5))

add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
             "Thanks for listening!", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")

add_text_box(slide, Inches(2.0), Inches(4.3), Inches(9), Inches(0.6),
             "Yoodle — Meetings that actually slap. \U0001f91f", font_size=24, color=YELLOW, bold=False, align=PP_ALIGN.CENTER, font_name="Arial")

# Links
add_shadow_rounded(slide, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.6), RGBColor(0x1A, 0x1A, 0x1A))
links = [
    {"text": "\U0001f310  yoodle.vercel.app", "size": 20, "color": CYAN, "bold": True, "space_after": 8, "align": PP_ALIGN.CENTER},
    {"text": "\U0001f4bb  github.com/uditanshutomar/Yoodle", "size": 18, "color": RGBColor(0xBB, 0xBB, 0xBB), "bold": False, "space_after": 8, "align": PP_ALIGN.CENTER},
    {"text": "Built with \u2764\ufe0f at HackCU 2025", "size": 16, "color": YELLOW, "bold": True, "space_after": 4, "align": PP_ALIGN.CENTER},
]
add_multi_text(slide, Inches(3.7), Inches(5.35), Inches(5.9), Inches(1.4), links, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────
output_path = "/Users/uditanshutomar/Desktop/Yoodle-HackCU-Pitch.pptx"
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
